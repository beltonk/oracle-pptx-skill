#!/usr/bin/env python3
"""Insert icons into PowerPoint slides based on content keywords.

Usage:
    python insert-icons.py <input.pptx> <output.pptx>
    
Analyzes slide content and inserts relevant icons from the icon library.
Converts SVG to PNG on-the-fly for PowerPoint compatibility.
"""

import sys
import json
import tempfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

def load_icon_index():
    """Load icon index from icon-index.json"""
    icon_index_path = Path(__file__).parent.parent / "resources" / "icons" / "icon-index.json"
    
    if not icon_index_path.exists():
        print("Warning: icon-index.json not found")
        return {}
    
    with open(icon_index_path, 'r') as f:
        return json.load(f)

# Keyword to icon mappings
ICON_KEYWORDS = {
    'ai': 'AI',
    'agent': 'Agents',
    'agentic': 'Agents',
    'autonomous': 'Autonomous',
    'database': 'Database',
    'cloud': 'Cloud',
    'finance': 'Financial',
    'operations': 'Operations',
    'customer': 'Customer',
    'supply': 'Supply',
    'architecture': 'Architecture',
    'security': 'Security',
    'governance': 'Governance',
    'business': 'Business',
    'enterprise': 'Enterprise',
    'data': 'Data',
    'analytics': 'Analytics',
    'automation': 'Automation',
    'strategy': 'Strategy',
    'experience': 'Experience',
}

def svg_to_png_placeholder(svg_path, output_path, size=300):
    """Convert SVG to PNG using available tools or create placeholder"""
    try:
        # Try using cairosvg if available
        import cairosvg
        cairosvg.svg2png(
            url=str(svg_path),
            write_to=str(output_path),
            output_width=size,
            output_height=size
        )
        return True
    except ImportError:
        pass
    
    try:
        # Try using PIL/Pillow with svglib
        from svglib.svglib import svg2rlg
        from reportlab.graphics import renderPM
        drawing = svg2rlg(svg_path)
        renderPM.drawToFile(drawing, output_path, fmt='PNG')
        return True
    except ImportError:
        pass
    
    # Fallback: create a simple placeholder PNG
    try:
        from PIL import Image, ImageDraw, ImageFont
        img = Image.new('RGBA', (size, size), (45, 49, 146, 255))  # Oracle blue
        draw = ImageDraw.Draw(img)
        
        # Draw simple icon placeholder
        margin = size // 4
        draw.rectangle([margin, margin, size-margin, size-margin], 
                      fill=(255, 255, 255, 200))
        
        img.save(output_path)
        return True
    except Exception as e:
        print(f"    Warning: Could not create icon placeholder: {e}")
        return False

def find_icon_by_keyword(keyword, icon_index):
    """Find best matching icon for a keyword"""
    keyword_lower = keyword.lower()
    
    # Try keyword mapping
    if keyword_lower in ICON_KEYWORDS:
        search_term = ICON_KEYWORDS[keyword_lower].lower()
        for category, icons in icon_index.get('categories', {}).items():
            for icon in icons:
                if search_term in icon['name'].lower():
                    return icon['file']
    
    # Try direct search
    for category, icons in icon_index.get('categories', {}).items():
        for icon in icons:
            if keyword_lower in icon['name'].lower():
                return icon['file']
    
    return None

def analyze_slide_content(slide):
    """Extract keywords from slide content"""
    text_parts = []
    
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame.text:
            text_parts.append(shape.text_frame.text)
    
    combined_text = ' '.join(text_parts).lower()
    
    # Find matching keywords
    matches = []
    for keyword in ICON_KEYWORDS.keys():
        if keyword in combined_text:
            matches.append((keyword, combined_text.count(keyword)))
    
    matches.sort(key=lambda x: x[1], reverse=True)
    return [m[0] for m in matches[:1]]  # Top keyword only

def insert_icon(slide, icon_png_path):
    """Insert an icon PNG into the slide"""
    try:
        # Position in top-right corner
        pic = slide.shapes.add_picture(
            str(icon_png_path),
            Inches(10.5), Inches(1.2),
            width=Inches(1.2), height=Inches(1.2)
        )
        return True
    except Exception as e:
        print(f"    Warning: Could not insert icon: {e}")
        return False

def process_presentation(input_pptx, output_pptx):
    """Process presentation and insert icons based on content"""
    
    icon_index = load_icon_index()
    if not icon_index:
        print("No icon index available")
        import shutil
        shutil.copy(input_pptx, output_pptx)
        return
    
    prs = Presentation(input_pptx)
    icons_dir = Path(__file__).parent.parent / "resources" / "icons" / "dark-theme"
    
    icons_inserted = 0
    
    # Create temp directory for PNG conversions
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        
        # Process each content slide
        for slide_idx, slide in enumerate(prs.slides):
            if slide_idx == 0 or slide_idx == len(prs.slides) - 1:
                continue  # Skip cover and thank you
            
            keywords = analyze_slide_content(slide)
            
            if keywords:
                keyword = keywords[0]
                icon_file = find_icon_by_keyword(keyword, icon_index)
                
                if icon_file:
                    svg_path = icons_dir / icon_file
                    if svg_path.exists():
                        # Convert to PNG
                        png_path = temp_path / f"icon_{slide_idx}.png"
                        if svg_to_png_placeholder(svg_path, png_path):
                            if insert_icon(slide, png_path):
                                icons_inserted += 1
                                print(f"  ✓ Slide {slide_idx + 1}: {keyword}")
        
        prs.save(output_pptx)
    
    print(f"\n✅ Icon insertion complete:")
    print(f"   Icons inserted: {icons_inserted}")
    print(f"   Saved to: {output_pptx}")

def main():
    if len(sys.argv) != 3:
        print(__doc__)
        sys.exit(1)
    
    input_pptx = Path(sys.argv[1])
    output_pptx = Path(sys.argv[2])
    
    if not input_pptx.exists():
        print(f"Error: Input file '{input_pptx}' not found")
        sys.exit(1)
    
    print(f"🎨 Inserting icons with keyword matching...\n")
    
    try:
        process_presentation(str(input_pptx), str(output_pptx))
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)

if __name__ == "__main__":
    main()
